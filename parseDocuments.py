import json
import os
import pickle
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

def parse_docx(file_path):
    doc = Document(file_path)
    text = '\n'.join([para.text for para in doc.paragraphs if para.text.strip()])
    return text

def parse_xlsx(file_path):
    wb = load_workbook(file_path, data_only=True)
    text = []
    for sheet in wb.worksheets:
        text.append(f"Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            row_text = '\t'.join([str(cell) if cell is not None else '' for cell in row])
            if row_text.strip():
                text.append(row_text)
    return '\n'.join(text)

def parse_pptx(file_path):
    prs = Presentation(file_path)
    text = []
    for i, slide in enumerate(prs.slides, 1):
        text.append(f"Slide {i}:")
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                text.append(shape.text)
    return '\n'.join(text)

def parse_document(file_path, file_type):
    try:
        if file_path.endswith('.docx'):
            return parse_docx(file_path)
        elif file_path.endswith('.xlsx'):
            return parse_xlsx(file_path)
        elif file_path.endswith('.pptx'):
            return parse_pptx(file_path)
        else:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
    except Exception as e:
        print(f"Error parsing {file_path}: {e}")
        return ""

def main():
    with open('assets/drive_files_metadata.json', 'r') as f:
        metadata = json.load(f)
    
    parsed_documents = []
    for item in metadata:
        if not os.path.exists(item['local_path']):
            continue
        
        content = parse_document(item['local_path'], item['file_type'])
        if content:
            parsed_documents.append({
                'page_id': item['page_id'],
                'page_title': item['page_title'],
                'source_url': item['source_url'],
                'file_type': item['file_type'],
                'content': content,
                'metadata': {
                    'title': f"{item['page_title']} - {os.path.basename(item['local_path'])}",
                    'id': item['page_id'],
                    'source': item['source_url'],
                    'file_path': item['local_path']
                }
            })
            print(f"Parsed: {item['page_title']}")
    
    with open('drive_documents_parsed.pkl', 'wb') as f:
        pickle.dump(parsed_documents, f)
    
    print(f"\nParsed {len(parsed_documents)} documents")

if __name__ == '__main__':
    main()
